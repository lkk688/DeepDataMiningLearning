import os
#pip install python-pptx Pillow
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, UnidentifiedImageError
import io
from pathlib import Path

def extract_from_pptx(pptx_path, output_dir):
    prs = Presentation(pptx_path)
    file_title = Path(pptx_path).stem
    markdown_content = f"# {file_title}\n\n"

    media_output_dir = os.path.join(output_dir, 'media')
    os.makedirs(media_output_dir, exist_ok=True)

    for idx, slide in enumerate(prs.slides):
        markdown_content += f"## Slide {idx + 1}\n\n"

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    markdown_content += f"{text}\n\n"

            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    image_bytes = image.blob
                    image_format = image.ext
                    
                    # Verify the image format is valid before saving
                    try:
                        # Try to open the image with PIL to validate it
                        with io.BytesIO(image_bytes) as image_file:
                            img = Image.open(image_file)
                            img.verify()  # Verify it's a valid image
                        
                        # If verification passes, save the image
                        image_filename = f"{file_title}_slide{idx + 1}_{shape.shape_id}.{image_format}"
                        image_path = os.path.join(media_output_dir, image_filename)

                        with open(image_path, 'wb') as f:
                            f.write(image_bytes)

                        markdown_content += f"![{image_filename}](media/{image_filename})\n\n"
                    
                    except UnidentifiedImageError:
                        print(f"Unsupported image format in {file_title}, slide {idx + 1}, shape ID {shape.shape_id}")
                        continue
                    except Exception as e:
                        print(f"Error processing image in {file_title}, slide {idx + 1}, shape ID {shape.shape_id}: {str(e)}")
                        continue
                        
                except Exception as e:
                    print(f"Error accessing image in {file_title}, slide {idx + 1}: {str(e)}")
                    continue

    # Write to markdown file
    output_md_path = os.path.join(output_dir, f"{file_title}.md")
    with open(output_md_path, 'w', encoding='utf-8') as f:
        f.write(markdown_content)

    print(f"Processed: {pptx_path}")

def process_all_pptx(input_folder, output_folder):
    os.makedirs(output_folder, exist_ok=True)
    for file in os.listdir(input_folder):
        if file.endswith('.pptx'):
            full_path = os.path.join(input_folder, file)
            try:
                extract_from_pptx(full_path, output_folder)
            except Exception as e:
                print(f"Failed to process {file}: {str(e)}")

# Example usage:
input_folder = "/Users/kaikailiu/Library/CloudStorage/GoogleDrive-kaikai.liu@sjsu.edu/My Drive/CurrentWork/Spring2025Courses/CMPE255-01 Spring2025/Final Slides"  # Update this path
output_folder = "/Users/kaikailiu/Documents/markdown_output"       # Update this path
process_all_pptx(input_folder, output_folder)